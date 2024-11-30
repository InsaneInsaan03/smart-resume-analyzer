from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def apply_text_style(shape):
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    # Style title
    title.text = "Smart Resume Analyzer"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Style subtitle
    subtitle.text = "A Mobile-First Resume Analysis and Career Guidance Application"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Style title
    title.text = "Overview"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content with proper margins
    content.text = """• Intelligent Resume Analysis System
• Mobile-First Approach
• Career Guidance Platform
• Skill Recommendation Engine
• Course Suggestion System"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.level = 0

def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # Add and style title
    title_shape = shapes.title
    title_shape.text = "System Architecture"
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Calculate center positions
    center_y = (prs.slide_height - Inches(2)) / 2
    total_width = Inches(11)  # Total width of all components
    start_x = (prs.slide_width - total_width) / 2
    
    # Create architecture diagram using shapes
    # User box
    user = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, center_y, Inches(2.5), Inches(1.5))
    user.text = "User/Mobile App"
    apply_text_style(user)
    
    # Arrow 1
    arrow1 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, start_x + Inches(2.5), center_y + Inches(0.25), Inches(1.5), Inches(1))
    
    # WebView box
    webview = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x + Inches(4), center_y, Inches(2.5), Inches(1.5))
    webview.text = "WebView Interface"
    apply_text_style(webview)
    
    # Arrow 2
    arrow2 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, start_x + Inches(6.5), center_y + Inches(0.25), Inches(1.5), Inches(1))
    
    # Streamlit box
    streamlit = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x + Inches(8), center_y, Inches(2.5), Inches(1.5))
    streamlit.text = "Streamlit Backend"
    apply_text_style(streamlit)
    
    # Style all shape text
    for shape in [user, webview, streamlit]:
        shape.text_frame.paragraphs[0].font.size = Pt(24)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Style title
    title.text = "Key Features"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content with proper formatting
    content.text = """Resume Analysis:
• PDF parsing and text extraction
• Skills identification
• Experience level assessment
• Career field detection

Career Guidance:
• Personalized skill recommendations
• Field-specific course suggestions
• Learning resource recommendations
• Career path guidance"""
    
    # Style paragraphs
    for paragraph in content.text_frame.paragraphs:
        if ":" in paragraph.text:  # Headers
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.level = 0
        else:  # Bullet points
            paragraph.font.size = Pt(24)
            paragraph.level = 1

def create_mobile_ui_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Style title
    title.text = "Mobile Interface"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content with proper formatting
    content.text = """Android App Features:
• Native Android WebView
• File upload capability
• Progress indicators
• Offline support
• Error handling
• Network state management"""
    
    # Style paragraphs
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        if i == 0:  # Header
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.level = 0
        else:  # Bullet points
            paragraph.font.size = Pt(24)
            paragraph.level = 1

def create_workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # Add and style title
    title = shapes.title
    title.text = "Application Workflow"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Create workflow diagram
    steps = [
        "Upload Resume",
        "Parse PDF",
        "Analyze Content",
        "Generate Recommendations",
        "Display Results"
    ]
    
    # Calculate center positions
    start_y = Inches(2)
    box_width = Inches(3)
    box_height = Inches(0.75)
    arrow_height = Inches(0.75)
    total_height = len(steps) * (box_height + arrow_height) - arrow_height
    start_x = (prs.slide_width - box_width) / 2
    
    for i, step in enumerate(steps):
        # Add box
        box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x,
            start_y + i * (box_height + arrow_height),
            box_width,
            box_height
        )
        
        # Style box
        tf = box.text_frame
        tf.text = step
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        apply_text_style(box)
        
        # Add arrow if not last step
        if i < len(steps) - 1:
            arrow = shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                start_x + (box_width - Inches(0.5)) / 2,
                start_y + box_height + i * (box_height + arrow_height),
                Inches(0.5),
                arrow_height
            )

def create_future_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Style title
    title.text = "Future Enhancements"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content with proper formatting
    content.text = """• Enhanced Analytics with AI/ML
• Resume Templates
• Interview Preparation Module
• Job Matching System
• iOS Support
• Desktop Application
• Cloud Synchronization
• Real-time Collaboration"""
    
    # Style paragraphs
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.level = 0

def create_presentation():
    prs = Presentation()
    
    # Set slide size to widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    create_title_slide(prs)
    create_overview_slide(prs)
    create_architecture_slide(prs)
    create_features_slide(prs)
    create_mobile_ui_slide(prs)
    create_workflow_slide(prs)
    create_future_slide(prs)
    
    # Save presentation with a different name
    prs.save('Resume_Analyzer_Slides.pptx')

if __name__ == "__main__":
    create_presentation()
